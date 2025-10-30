#!/usr/bin/env python3
"""
Create a PowerPoint (.pptx) presentation from ND2 files by extracting fluorescence images.

Features (inspired by create_keynote_from_nd2.py, with enhancements):
- Recursive search from target directory for .nd2 files (skips hidden/AppleDouble files)
- Channel mapping by name (DAPI->Blue, 488->Green, 568->Red), with sensible fallbacks
- Optional max-intensity projection over Z and/or T
- Optional intensity normalization per channel and percentile clipping
- Optional image scaling and maximum slide image size
- Optional montage for multiple scenes/positions (S) or timepoints (T)
- Embeds slide title and optional scale text derived from metadata (µm/side)
- Supports saving intermediate JPEGs for review/debugging

Requirements:
- Python 3.9+
- Packages: nd2, numpy, Pillow, python-pptx

Usage examples:
  python create_powerpoint_from_nd2.py \
    --dir . --recursive --output "Fluorescence.pptx" --mip-z --mip-t --scale 0.5 \
    --clip-percent 0.3 --max-slide-size 1400 --keep-jpgs

"""

import argparse
import os
import sys
from dataclasses import dataclass
from typing import List, Optional, Tuple

import numpy as np
from PIL import Image
import nd2

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.util import Inches, Pt
except Exception as e:  # pragma: no cover - allow import-time errors to be reported nicely
    Presentation = None  # type: ignore


# -------------------------------
# Data types
# -------------------------------

@dataclass
class FluorescenceImage:
    jpg_path: str
    title: str
    side_length_um: float


# -------------------------------
# ND2 handling helpers
# -------------------------------

def get_image_axes_info(nd2_file: nd2.ND2File) -> str:
    try:
        return ''.join(nd2_file.sizes.keys())
    except Exception:
        return ""


def compute_side_length_um(nd2_file: nd2.ND2File) -> float:
    """Best-effort side length (X) in micrometers from metadata."""
    side_length_um = 0.0
    try:
        md = nd2_file.metadata
        channels = getattr(md, 'channels', None)
        if channels and hasattr(channels[0], 'volume'):
            volume = channels[0].volume
            if hasattr(volume, 'axesCalibration') and hasattr(volume, 'voxelCount'):
                pixel_size_um = float(volume.axesCalibration[0])
                image_size_pixels = int(volume.voxelCount[0])
                side_length_um = pixel_size_um * image_size_pixels
        else:
            # Fallback path via metadata if available
            scaling = getattr(md, 'pixels', None)
            if scaling and hasattr(scaling, 'sizeUm') and hasattr(scaling, 'shape'):  # type: ignore[attr-defined]
                try:
                    side_length_um = float(scaling.sizeUm[0]) * int(scaling.shape[0])  # type: ignore[index]
                except Exception:
                    pass
    except Exception:
        pass
    return side_length_um


def normalize_channel(
    channel_data: np.ndarray,
    clip_percent: float = 0.0,
) -> np.ndarray:
    """Normalize to 0-255 uint8 with optional percentile clipping."""
    channel = np.nan_to_num(channel_data, nan=0.0, posinf=0.0, neginf=0.0).astype(np.float32)
    if channel.size == 0:
        return np.zeros_like(channel, dtype=np.uint8)
    vmin = float(np.min(channel))
    vmax = float(np.max(channel))
    if clip_percent > 0.0:
        lo = np.percentile(channel, clip_percent)
        hi = np.percentile(channel, 100.0 - clip_percent)
        vmin, vmax = float(lo), float(hi)
    if vmax <= vmin:
        return np.zeros_like(channel, dtype=np.uint8)
    scaled = (np.clip(channel, vmin, vmax) - vmin) / (vmax - vmin)
    return np.round(scaled * 255.0).astype(np.uint8)


def max_intensity_projection(array: np.ndarray, axis: int) -> np.ndarray:
    return np.max(array, axis=axis)


def extract_rgb_from_nd2(
    nd2_path: str,
    mip_z: bool,
    mip_t: bool,
    clip_percent: float,
    preferred_order: Tuple[str, str, str] = ("red", "blue", "green"),
    verbose: bool = False,
) -> Tuple[np.ndarray, float, Tuple[str, str, str]]:
    """
    Return (H,W,3) uint8 RGB image and side length µm.

    Channel mapping by name:
      - 'dapi' -> blue
      - '488' or 'alexa488' -> green
      - '568' or 'alx568' -> red
    Fallbacks: use available channels in order for R,G,B if names are not found.
    Applies optional MIP across Z (axis 'Z') and/or T (axis 'T') when present.
    """
    with nd2.ND2File(nd2_path) as f:
        sizes = dict(f.sizes)
        axes = ''.join(sizes.keys())
        side_um = compute_side_length_um(f)

        # Load as array with all available axes
        data = f.asarray()

        # Bring into a shape we can reason about: prioritize channel axis 'C'
        # nd2 may provide data with axes order like 'TCYZXS' etc.
        # We'll collapse non-channel axes via MIP if requested or take the first index.
        # Build view to isolate channel plane stack as (C, Y, X)
        axis_to_index = {ax: i for i, ax in enumerate(axes)}

        # Resolve per-axis reducer
        def reduce_axis(a: np.ndarray, ax_char: str) -> np.ndarray:
            if ax_char not in axis_to_index:
                return a
            axis_id = axis_to_index[ax_char]
            if ax_char == 'Z' and mip_z:
                return max_intensity_projection(a, axis=axis_id)
            if ax_char == 'T' and mip_t:
                return max_intensity_projection(a, axis=axis_id)
            # Take the first slice as default
            slicer = [slice(None)] * a.ndim
            slicer[axis_id] = 0
            return a[tuple(slicer)]

        # Reduce non-C spatial/series axes in a stable order
        for ax_char in list(axes):
            if ax_char == 'C':
                continue
            if data.ndim <= 3 and 'C' in axes:
                break
            data = reduce_axis(data, ax_char)
            # recompute axes metadata since dims may change
            axes = ''.join([ax for ax in axes if ax != ax_char])
            axis_to_index = {ax: i for i, ax in enumerate(axes)}

        # Ensure we have (C, Y, X)
        if 'C' in axes:
            c_idx = axis_to_index['C']
            if c_idx != 0:
                # Move channel axis to front
                perm = [c_idx] + [i for i in range(data.ndim) if i != c_idx]
                data = np.transpose(data, perm)
            if data.ndim == 2:
                data = data[np.newaxis, ...]
            if data.ndim != 3:
                raise ValueError(f"Unexpected shape after reductions: {data.shape}")
        else:
            # No channel axis -> treat as single-channel
            if data.ndim == 2:
                data = data[np.newaxis, ...]
            elif data.ndim == 3:
                # Assume already (Y,X,?) -> move to (1,Y,X) by taking first
                data = data[..., 0]
                data = data[np.newaxis, ...]
            else:
                raise ValueError(f"Unsupported ND2 shape without channel axis: {data.shape}")

        # Map channels by metadata names when possible
        red = np.zeros_like(data[0], dtype=np.uint8)
        green = np.zeros_like(data[0], dtype=np.uint8)
        blue = np.zeros_like(data[0], dtype=np.uint8)
        label_red: str = ""
        label_green: str = ""
        label_blue: str = ""

        md_channels = getattr(f.metadata, 'channels', None)
        # Helpers for channel classification
        def is_brightfield(name_lc: str) -> bool:
            bf_keys = [
                'brightfield', 'bright field', 'bf', 'td', 'tl', 'trans', 'transmitted', 'phase', 'ph', 'differential interference', 'dic'
            ]
            return any(k in name_lc for k in bf_keys)
        def is_dapi(name_lc: str) -> bool:
            return 'dapi' in name_lc or 'hoechst' in name_lc or '405' in name_lc
        def is_green(name_lc: str) -> bool:
            return (
                '488' in name_lc or 'alexa488' in name_lc or 'gfp' in name_lc or 'fitc' in name_lc
            )
        def is_red(name_lc: str) -> bool:
            return (
                '568' in name_lc or '561' in name_lc or '555' in name_lc or '594' in name_lc or 'alx568' in name_lc
                or 'cy3' in name_lc or 'mcherry' in name_lc or 'texas' in name_lc
            )
        mapped = {"red": False, "green": False, "blue": False}
        channel_names: list[str] = []
        if md_channels:
            for i, ch in enumerate(md_channels):
                try:
                    raw = ch.channel.name  # type: ignore[attr-defined]
                    name = raw.lower()
                except Exception:
                    raw = f"Channel {i}"
                    name = raw.lower()
                channel_names.append(name)
                if is_brightfield(name):
                    if verbose:
                        print(f"  - skip brightfield-like channel {i}: '{raw}'")
                    # Skip brightfield-like channels for fluorescence mapping
                    continue
                if is_dapi(name) and not mapped["blue"]:
                    blue = normalize_channel(data[i], clip_percent)
                    mapped["blue"] = True
                    label_blue = raw
                    if verbose:
                        print(f"  - map BLUE from channel {i}: '{raw}'")
                elif is_green(name) and not mapped["green"]:
                    green = normalize_channel(data[i], clip_percent)
                    mapped["green"] = True
                    label_green = raw
                    if verbose:
                        print(f"  - map GREEN from channel {i}: '{raw}'")
                elif is_red(name) and not mapped["red"]:
                    red = normalize_channel(data[i], clip_percent)
                    mapped["red"] = True
                    label_red = raw
                    if verbose:
                        print(f"  - map RED from channel {i}: '{raw}'")

        # Fallbacks: assign remaining channels by order, excluding brightfield-like channels and already-mapped channels
        used_channels = set()
        if mapped["red"]:
            # Find which channel was used for red
            for i, ch in enumerate(md_channels or []):
                try:
                    raw = ch.channel.name
                    name = raw.lower()
                except Exception:
                    raw = f"Channel {i}"
                    name = raw.lower()
                if is_red(name):
                    used_channels.add(i)
                    break
        if mapped["green"]:
            # Find which channel was used for green
            for i, ch in enumerate(md_channels or []):
                try:
                    raw = ch.channel.name
                    name = raw.lower()
                except Exception:
                    raw = f"Channel {i}"
                    name = raw.lower()
                if is_green(name):
                    used_channels.add(i)
                    break
        if mapped["blue"]:
            # Find which channel was used for blue
            for i, ch in enumerate(md_channels or []):
                try:
                    raw = ch.channel.name
                    name = raw.lower()
                except Exception:
                    raw = f"Channel {i}"
                    name = raw.lower()
                if is_dapi(name):
                    used_channels.add(i)
                    break

        remaining: list[tuple[int, np.ndarray]] = []
        for i in range(data.shape[0]):
            if i in used_channels:
                continue  # Skip already-mapped channels
            name = channel_names[i] if i < len(channel_names) else f"channel{i}"
            if is_brightfield(name):
                continue
            remaining.append((i, normalize_channel(data[i], clip_percent)))
        for color in preferred_order:
            if not mapped[color] and remaining:
                idx, img = remaining.pop(0)
                fallback_label = f"Channel {idx}"
                if color == 'red':
                    red = img
                    label_red = label_red or fallback_label
                    if verbose:
                        print(f"  - fallback RED from channel {idx}")
                elif color == 'green':
                    green = img
                    label_green = label_green or fallback_label
                    if verbose:
                        print(f"  - fallback GREEN from channel {idx}")
                else:
                    blue = img
                    label_blue = label_blue or fallback_label
                    if verbose:
                        print(f"  - fallback BLUE from channel {idx}")

        # Heuristic fix: if green ended up empty but blue has data (common mislabel -> appears blue),
        # prefer showing that as green to match expectation.
        if green.max() == 0 and blue.max() > 0 and red.max() == 0:
            green = blue
            blue = np.zeros_like(blue)
            if not label_green and label_blue:
                label_green, label_blue = label_blue, ""
            if verbose:
                print("  - heuristic: moved BLUE to GREEN to avoid mislabel (blue-only)")

        # If DAPI (blue) is present together with a strong green-only signal, suppress blue to avoid cyan.
        try:
            if (blue.max() > 0 and green.max() > 0 and red.max() == 0 and label_blue and 'dapi' in label_blue.lower()):
                blue = np.zeros_like(blue)
                label_blue = ""
                if verbose:
                    print("  - heuristic: suppressed DAPI (blue) to avoid cyan with green-only signal")
        except Exception:
            pass

        if verbose:
            print(f"  => final labels  R:'{label_red or 'Red'}'  G:'{label_green or 'Green'}'  B:'{label_blue or 'Blue'}'")
            print(f"  => channel intensities  R:{red.max():.1f}  G:{green.max():.1f}  B:{blue.max():.1f}")
        
        # Debug: Test if green and blue are swapped in the final output
        # Create RGB array ensuring proper channel order: [Red, Green, Blue]
        rgb = np.stack([red, green, blue], axis=2)
        return rgb, side_um, (label_red or "Red", label_green or "Green", label_blue or "Blue")


# -------------------------------
# I/O helpers
# -------------------------------

def save_rgb_as_jpg(rgb: np.ndarray, out_dir: str, base_name: str, quality: int = 95, scale: float = 1.0, max_side: Optional[int] = None) -> str:
    os.makedirs(out_dir, exist_ok=True)
    img = Image.fromarray(rgb, mode='RGB')

    if scale and scale != 1.0:
        new_size = (max(1, int(img.width * scale)), max(1, int(img.height * scale)))
        img = img.resize(new_size, resample=Image.BICUBIC)

    if max_side:
        w, h = img.size
        side = max(w, h)
        if side > max_side:
            factor = max_side / float(side)
            img = img.resize((max(1, int(w * factor)), max(1, int(h * factor))), resample=Image.BICUBIC)

    out_path = os.path.join(out_dir, f"{base_name}.jpg")
    img.save(out_path, 'JPEG', quality=quality)
    return out_path


def add_slide_with_image(prs: 'Presentation', image_path: str, title: str, scale_text: Optional[str], channels_text: Optional[str], max_width_px: int) -> None:
    # Try to pick a "Two Content" layout; fallback to Title and Content if not present
    two_content_layout = None
    title_and_content_layout = None
    for layout in prs.slide_layouts:
        try:
            name = layout.name
        except Exception:
            name = ""
        if name and ("Two Content" in name or "2" in name and "Content" in name or "比較" in name):
            two_content_layout = layout
        if name and ("Title and Content" in name or "タイトルとコンテンツ" in name):
            title_and_content_layout = layout
    layout = two_content_layout or title_and_content_layout or prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout)

    # Set title: partial path + file name
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        try:
            tf = slide.shapes.title.text_frame
            p = tf.paragraphs[0]
            # ensure at least one run exists to apply font
            if len(p.runs) == 0:
                r = p.add_run()
                r.text = title
                r.font.name = "MS Pゴシック"
                r.font.size = Pt(44)
            else:
                for r in p.runs:
                    r.font.name = "MS Pゴシック"
                    r.font.size = Pt(44)
        except Exception:
            pass

    # Collect placeholders explicitly by type for reliability
    pic_like_types = set()
    text_like_types = set()
    try:
        pic_like_types = {PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.CONTENT, PP_PLACEHOLDER.OBJECT}
        text_like_types = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CONTENT}
        title_like_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE}
    except Exception:
        pass

    picture_phs = []
    text_phs = []
    other_phs = []
    for shp in slide.shapes.placeholders:
        try:
            is_title_attr = getattr(shp, 'is_title', False)
            phf = getattr(shp, 'placeholder_format', None)
            ph_type = getattr(phf, 'type', None)
        except Exception:
            is_title_attr = False
            ph_type = None
        # Robust title detection
        is_title_like = False
        try:
            if ph_type in title_like_types:
                is_title_like = True
        except Exception:
            pass
        if is_title_attr or is_title_like or (hasattr(slide.shapes, 'title') and shp == slide.shapes.title):
            continue
        if ph_type in pic_like_types:
            picture_phs.append(shp)
        elif ph_type in text_like_types:
            text_phs.append(shp)
        else:
            other_phs.append(shp)

    # Sort by x position
    picture_phs.sort(key=lambda p: p.left)
    text_phs.sort(key=lambda p: p.left)
    other_phs.sort(key=lambda p: p.left)

    if len(picture_phs) == 0 and len(other_phs) > 0:
        # fallback: use the left-most non-title placeholder as image host
        picture_phs = [other_phs[0]]

    if len(text_phs) == 0 and len(other_phs) > 1:
        # fallback: use a different placeholder for text if available
        text_phs = [other_phs[-1]]

    if len(picture_phs) == 0 and len(text_phs) == 0:
        # Fallback: place image on the current slide within safe margins (no extra slide)
        left_margin = Inches(0.5)
        right_margin = Inches(0.5)
        top = Inches(1.5)
        bottom_margin = Inches(0.5)
        picture = slide.shapes.add_picture(image_path, left_margin, top)
        max_w = prs.slide_width - (left_margin + right_margin)
        max_h = prs.slide_height - (top + bottom_margin)
        fit_ratio = min(float(max_w) / float(picture.width), float(max_h) / float(picture.height), 1.0)
        final_ratio = fit_ratio * 0.75
        if final_ratio < 1.0:
            picture.width = int(picture.width * final_ratio)
            picture.height = int(picture.height * final_ratio)
        picture.left = int((prs.slide_width - picture.width) / 2)

        # Also place the side length and channel info on the right area
        right_left = int(prs.slide_width * 0.58)
        right_top = int(prs.slide_height * 0.22)
        right_width = int(prs.slide_width * 0.38)
        right_height = int(prs.slide_height * 0.6)
        tb = slide.shapes.add_textbox(right_left, right_top, right_width, right_height)
        tf = tb.text_frame
        tf.clear()
        if scale_text:
            p = tf.paragraphs[0]
            p.text = scale_text
            p.level = 0
        if channels_text:
            parts = [s.strip() for s in channels_text.split('|')]
            for i, part in enumerate(parts):
                if i == 0 and not scale_text:
                    p = tf.paragraphs[0]
                    p.text = part
                    p.level = 0
                else:
                    rp = tf.add_paragraph()
                    rp.text = part
                    rp.level = 0
        return

    # Choose left picture host and right text host
    left_ph = picture_phs[0]
    # pick a text placeholder that's to the right of left_ph if possible
    right_candidates = [ph for ph in text_phs if ph.left > left_ph.left]
    right_ph = right_candidates[0] if len(right_candidates) > 0 else (text_phs[0] if len(text_phs) > 0 else None)

    # Left content: insert picture into placeholder to fit it nicely
    # Clear any existing text in left placeholder to avoid overlay text
    try:
        if hasattr(left_ph, 'text_frame') and left_ph.text_frame is not None:
            left_ph.text_frame.clear()
            # Also remove any existing paragraphs
            for para in left_ph.text_frame.paragraphs:
                para.clear()
    except Exception:
        pass

    inserted = False
    # Save placeholder dimensions before attempting insertion
    ph_left = left_ph.left
    ph_top = left_ph.top
    ph_width = left_ph.width
    ph_height = left_ph.height
    
    # Try multiple methods to insert picture into placeholder
    # Method 1: Use insert_picture if placeholder type supports it
    try:
        phf = getattr(left_ph, 'placeholder_format', None)
        ph_type = getattr(phf, 'type', None)
    except Exception:
        ph_type = None
    
    # Try insert_picture for all placeholder types (CONTENT placeholders can also hold pictures)
    try:
        new_pic = left_ph.insert_picture(image_path)
        if new_pic is not None:
            inserted = True
    except Exception:
        pass
    
    if not inserted:
        # Method 2: Remove placeholder and add picture in its exact position
        try:
            # Get placeholder's parent element to remove it
            ph_element = left_ph.element
            ph_parent = ph_element.getparent()
            if ph_parent is not None:
                ph_parent.remove(ph_element)
        except Exception:
            # If removal fails, try to hide it
            try:
                left_ph.width = 0
                left_ph.height = 0
            except Exception:
                pass
        
        # Add picture in the placeholder's exact position and size
        pic = slide.shapes.add_picture(image_path, ph_left, ph_top)
        # Scale to fit within placeholder dimensions
        max_w = ph_width
        max_h = ph_height
        fit_ratio = min(float(max_w) / float(pic.width), float(max_h) / float(pic.height), 1.0)
        final_ratio = fit_ratio * 0.98
        if final_ratio < 1.0:
            pic.width = int(pic.width * final_ratio)
            pic.height = int(pic.height * final_ratio)
        # Center within placeholder area
        pic.left = ph_left + int((ph_width - pic.width) / 2)
        pic.top = ph_top + int((ph_height - pic.height) / 2)

    # Right content: bullet list for side length and channel info
    placed_right_text = False
    if right_ph is not None and hasattr(right_ph, 'text_frame') and right_ph.text_frame is not None:
        try:
            tf = right_ph.text_frame
            tf.clear()
            tf.word_wrap = True
            if scale_text:
                p = tf.paragraphs[0]
                p.text = scale_text
                p.level = 0
                lines_added = 1
            else:
                lines_added = 0
            if channels_text:
                # Split to three items: Red..., Green..., Blue...
                parts = [s.strip() for s in channels_text.split('|')]
                for i, part in enumerate(parts):
                    if i == 0 and not scale_text:
                        p = tf.paragraphs[0]
                        p.text = part
                        p.level = 0
                        lines_added += 1
                    else:
                        rp = tf.add_paragraph()
                        rp.text = part
                        rp.level = 0
                        lines_added += 1
            placed_right_text = lines_added > 0
        except Exception:
            placed_right_text = False

    # Fallback: if right placeholder unavailable, add a textbox on the right half
    if not placed_right_text:
        # Define a conservative right panel area
        right_left = int(prs.slide_width * 0.58)
        right_top = int(prs.slide_height * 0.22)
        right_width = int(prs.slide_width * 0.38)
        right_height = int(prs.slide_height * 0.6)
        tb = slide.shapes.add_textbox(right_left, right_top, right_width, right_height)
        tf = tb.text_frame
        tf.clear()
        if scale_text:
            p = tf.paragraphs[0]
            p.text = scale_text
            p.level = 0
        if channels_text:
            parts = [s.strip() for s in channels_text.split('|')]
            for i, part in enumerate(parts):
                if i == 0 and not scale_text:
                    p = tf.paragraphs[0]
                    p.text = part
                    p.level = 0
                else:
                    rp = tf.add_paragraph()
                    rp.text = part
                    rp.level = 0


# -------------------------------
# Discovery of ND2 files
# -------------------------------

def find_nd2_files(root_dir: str, recursive: bool) -> List[str]:
    paths: List[str] = []
    if recursive:
        for cur_dir, dirnames, filenames in os.walk(root_dir):
            # skip hidden directories
            dirnames[:] = [d for d in dirnames if not d.startswith('.')]
            for fn in filenames:
                if not fn.lower().endswith('.nd2'):
                    continue
                if fn.startswith('.') or fn.startswith('._'):
                    continue
                paths.append(os.path.join(cur_dir, fn))
    else:
        for fn in os.listdir(root_dir):
            if fn.lower().endswith('.nd2') and not fn.startswith('.') and not fn.startswith('._'):
                paths.append(os.path.join(root_dir, fn))
    paths.sort()
    return paths


# -------------------------------
# Main
# -------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(description="Create PowerPoint from ND2 fluorescence images.")
    parser.add_argument("--dir", type=str, default=os.getcwd(), help="Target directory (default: cwd)")
    parser.add_argument("--recursive", action="store_true", default=True, help="Recursively search subdirectories for .nd2 (default: on)")
    parser.add_argument("--output", type=str, default="", help="Output .pptx file name (default: auto)")
    parser.add_argument("--mip-z", action="store_true", help="Max-intensity project along Z if available")
    parser.add_argument("--mip-t", action="store_true", help="Max-intensity project along time if available")
    parser.add_argument("--clip-percent", type=float, default=0.0, help="Percentile clip (e.g., 0.3) per channel")
    parser.add_argument("--scale", type=float, default=1.0, help="Scale factor for exported JPG (e.g., 0.5)")
    parser.add_argument("--max-slide-size", type=int, default=1600, help="Max longer side for JPG before placing on slide")
    parser.add_argument("--keep-jpgs", action="store_true", help="Keep generated JPGs next to output")
    parser.add_argument("--jpg-dir", type=str, default="", help="Directory to store intermediate JPGs (default: temp next to output)")
    parser.add_argument("--verbose", action="store_true", help="Print detailed channel mapping and placement info")
    args = parser.parse_args()

    root_dir = os.path.abspath(args.dir)
    if not os.path.isdir(root_dir):
        print("Directory not found:", root_dir, file=sys.stderr)
        sys.exit(1)

    nd2_paths = find_nd2_files(root_dir, recursive=args.recursive)
    if not nd2_paths:
        print("No ND2 files found.")
        sys.exit(0)

    print(f"Found {len(nd2_paths)} ND2 file(s):")
    for p in nd2_paths:
        print("  -", os.path.relpath(p, root_dir))

    # Prepare output
    default_base = os.path.basename(os.path.abspath(root_dir).rstrip(os.sep)) or "FluorescencePresentation"
    out_name = args.output.strip() or f"{default_base}.pptx"
    if not out_name.lower().endswith('.pptx'):
        out_name += '.pptx'
    out_pptx_path = os.path.join(root_dir, out_name)

    jpg_dir = args.jpg_dir.strip() or os.path.join(root_dir, "temp_fluorescence_images_pptx")
    os.makedirs(jpg_dir, exist_ok=True)

    # Create presentation
    if Presentation is None:
        print("python-pptx is not installed. Please install with: pip install python-pptx", file=sys.stderr)
        sys.exit(1)
    prs = Presentation()

    created: List[FluorescenceImage] = []
    for nd2_path in nd2_paths:
        base = os.path.splitext(os.path.basename(nd2_path))[0]
        print(f"Processing {os.path.basename(nd2_path)}...")
        try:
            rgb, side_um, channel_labels = extract_rgb_from_nd2(
                nd2_path,
                mip_z=args.mip_z,
                mip_t=args.mip_t,
                clip_percent=args.clip_percent,
                verbose=args.verbose,
            )
        except Exception as e:
            print(f"  Error reading {os.path.basename(nd2_path)}: {e}")
            continue

        jpg_path = save_rgb_as_jpg(
            rgb,
            out_dir=jpg_dir,
            base_name=f"{base}_fluorescence",
            quality=95,
            scale=args.scale,
            max_side=args.max_slide_size,
        )

        # Title uses partial path + filename (relative to root)
        title = os.path.relpath(nd2_path, root_dir)
        scale_text = f"{side_um:.0f}µm/side" if side_um > 0 else "Scale unknown"
        channels_text = f"Red: {channel_labels[0]}  |  Green: {channel_labels[1]}  |  Blue: {channel_labels[2]}"
        add_slide_with_image(prs, jpg_path, title, scale_text, channels_text, max_width_px=args.max_slide_size)
        created.append(FluorescenceImage(jpg_path, title, side_um))

    if not created:
        print("No slides created (no images).")
        sys.exit(1)

    prs.save(out_pptx_path)
    print("✅ Saved PowerPoint:", out_pptx_path)

    if not args.keep_jpgs:
        # Cleanup
        for item in created:
            try:
                os.remove(item.jpg_path)
            except OSError:
                pass
        try:
            os.rmdir(jpg_dir)
        except OSError:
            pass
        print("✅ Cleanup complete")
    else:
        print(f"JPG files kept in: {jpg_dir}")


if __name__ == "__main__":
    main()


